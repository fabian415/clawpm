#!/usr/bin/env python3
"""Extract content images embedded in meeting supplementary documents (PPTX/XLSX/PDF)."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
from datetime import datetime, timezone
from pathlib import Path

DEFAULT_MIN_SIZE = 120
DEFAULT_MAX_IMAGES = 200
SUPPORTED_SUFFIXES = {".pptx", ".xlsx", ".pdf"}


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Extract embedded content images from PPTX/XLSX/PDF supplementary files."
    )
    parser.add_argument("inputs", type=Path, nargs="+", help="Supplementary document paths (.pptx/.xlsx/.pdf)")
    parser.add_argument("--output-dir", type=Path, required=True, help="Directory to write extracted images into")
    parser.add_argument("--manifest", type=Path, help="Destination manifest JSON path (default: <output-dir>/manifest.json)")
    parser.add_argument("--min-size", type=int, default=DEFAULT_MIN_SIZE,
                         help="Minimum width AND height in pixels to keep an image (filters icons/logos); default: %(default)s")
    parser.add_argument("--max-images", type=int, default=DEFAULT_MAX_IMAGES,
                         help="Safety cap on total images extracted across all inputs; default: %(default)s")
    args = parser.parse_args()

    args.output_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = (args.manifest or args.output_dir / "manifest.json").resolve()

    seen_hashes: set[str] = set()
    images: list[dict] = []
    skipped: list[str] = []

    for input_path in args.inputs:
        if not input_path.exists():
            skipped.append(f"{input_path}: 檔案不存在")
            continue
        suffix = input_path.suffix.lower()
        if suffix not in SUPPORTED_SUFFIXES:
            skipped.append(f"{input_path}: 不支援的格式（僅支援 .pptx/.xlsx/.pdf）")
            continue
        try:
            if suffix == ".pptx":
                found = extract_pptx(input_path, args.output_dir, args.min_size, seen_hashes)
            elif suffix == ".xlsx":
                found = extract_xlsx(input_path, args.output_dir, args.min_size, seen_hashes)
            else:
                found = extract_pdf(input_path, args.output_dir, args.min_size, seen_hashes)
        except RuntimeError as exc:
            skipped.append(f"{input_path}: {exc}")
            continue

        remaining = args.max_images - len(images)
        if remaining <= 0:
            skipped.append(f"{input_path}: 已達 --max-images 上限，略過剩餘圖片")
            break
        images.extend(found[:remaining])

    manifest = {
        "generatedAt": datetime.now(timezone.utc).isoformat(),
        "outputDir": str(args.output_dir.resolve()),
        "images": images,
        "skipped": skipped,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Extracted {len(images)} image(s) from {len(args.inputs)} file(s); manifest saved to {manifest_path}")
    for img in images:
        print(f"  {img['file']}  <-  {img['sourceFile']} ({img['location']}, {img['width']}x{img['height']})")
    for msg in skipped:
        print(f"  [skip] {msg}")


def _slugify(stem: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9_-]+", "-", stem).strip("-")
    return slug or "file"


def _unique_prefix(path: Path) -> str:
    digest = hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()[:8]
    return f"{_slugify(path.stem)}-{digest}"


def _save_image(output_dir: Path, prefix: str, location_tag: str, seq: int, ext: str,
                 data: bytes, seen_hashes: set[str]) -> str | None:
    digest = hashlib.sha1(data).hexdigest()
    if digest in seen_hashes:
        return None
    seen_hashes.add(digest)
    ext = ext.lstrip(".").lower() or "png"
    filename = f"{prefix}_{location_tag}_{seq:02d}.{ext}"
    (output_dir / filename).write_bytes(data)
    return filename


def extract_pptx(path: Path, output_dir: Path, min_size: int, seen_hashes: set[str]) -> list[dict]:
    """Render each slide as a full-page screenshot instead of pulling embedded
    media objects, since the latter mostly surfaces decorative icons/logos
    rather than the slide's actual visual content.

    The PPTX->PDF conversion is delegated to a remote `unoserver` sidecar
    (see docker/converter/Dockerfile) instead of a local LibreOffice install,
    so this image only needs the lightweight `unoserver` pip client."""
    import os
    import shutil
    import subprocess
    import tempfile

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to read PPTX slide captions (pip install python-pptx)") from exc
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("pymupdf is required to render PPTX slides (pip install pymupdf)") from exc

    unoconvert = shutil.which("unoconvert")
    if not unoconvert:
        raise RuntimeError("unoconvert is required to render PPTX slides (pip install unoserver)")

    host = os.environ.get("UNOSERVER_HOST", "host.docker.internal")
    port = os.environ.get("UNOSERVER_PORT", "2003")

    prefix = _unique_prefix(path)
    captions = [_pptx_slide_caption(slide) for slide in Presentation(path).slides]
    results = []

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        pdf_path = tmp_dir / f"{path.stem}.pdf"
        proc = subprocess.run(
            [unoconvert, "--host", host, "--port", str(port), "--host-location", "remote",
             "--convert-to", "pdf", str(path), str(pdf_path)],
            capture_output=True, text=True, timeout=180,
        )
        if proc.returncode != 0 or not pdf_path.exists():
            raise RuntimeError(
                f"unoserver ({host}:{port}) 轉換 PPTX 失敗，請確認 clawpm-unoserver 容器是否在跑: "
                f"{(proc.stderr or proc.stdout).strip()[:200]}"
            )

        doc = fitz.open(str(pdf_path))
        try:
            matrix = fitz.Matrix(2.0, 2.0)  # ~144 DPI，足夠辨識文字與圖表細節
            for slide_idx in range(len(doc)):
                pix = doc[slide_idx].get_pixmap(matrix=matrix)
                if pix.width < min_size or pix.height < min_size:
                    continue
                filename = _save_image(output_dir, prefix, f"slide{slide_idx + 1}", 1, "png",
                                        pix.tobytes("png"), seen_hashes)
                if not filename:
                    continue
                results.append({
                    "file": filename,
                    "sourceFile": path.name,
                    "sourceType": "pptx",
                    "location": f"投影片 {slide_idx + 1}",
                    "width": pix.width,
                    "height": pix.height,
                    "captionHint": captions[slide_idx] if slide_idx < len(captions) else "",
                })
        finally:
            doc.close()
    return results


def _pptx_slide_caption(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().splitlines()[0][:120]
    return ""


def extract_xlsx(path: Path, output_dir: Path, min_size: int, seen_hashes: set[str]) -> list[dict]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to read XLSX files (pip install openpyxl)") from exc

    prefix = _unique_prefix(path)
    wb = load_workbook(path, data_only=True)
    results = []

    for ws in wb.worksheets:
        sheet_images = getattr(ws, "_images", [])
        seq = 0
        for img in sheet_images:
            try:
                data = img._data()
            except Exception:
                continue
            width, height = int(getattr(img, "width", 0) or 0), int(getattr(img, "height", 0) or 0)
            if width < min_size or height < min_size:
                continue
            ext = (getattr(img, "format", None) or "png").lower()
            seq += 1
            sheet_tag = _slugify(ws.title)
            filename = _save_image(output_dir, prefix, f"{sheet_tag}", seq, ext, data, seen_hashes)
            if not filename:
                continue
            results.append({
                "file": filename,
                "sourceFile": path.name,
                "sourceType": "xlsx",
                "location": f"工作表「{ws.title}」",
                "width": width,
                "height": height,
                "captionHint": "",
            })
    return results


def extract_pdf(path: Path, output_dir: Path, min_size: int, seen_hashes: set[str]) -> list[dict]:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("pymupdf is required to read PDFs (pip install pymupdf)") from exc

    prefix = _unique_prefix(path)
    doc = fitz.open(str(path))
    results = []

    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            caption = (page.get_text().strip().splitlines() or [""])[0][:120]
            seq = 0
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                base = doc.extract_image(xref)
                width, height = base.get("width", 0), base.get("height", 0)
                if width < min_size or height < min_size:
                    continue
                seq += 1
                filename = _save_image(output_dir, prefix, f"p{page_index + 1}", seq, base.get("ext", "png"),
                                        base["image"], seen_hashes)
                if not filename:
                    continue
                results.append({
                    "file": filename,
                    "sourceFile": path.name,
                    "sourceType": "pdf",
                    "location": f"第 {page_index + 1} 頁",
                    "width": width,
                    "height": height,
                    "captionHint": caption,
                })
    finally:
        doc.close()
    return results


if __name__ == "__main__":
    main()
