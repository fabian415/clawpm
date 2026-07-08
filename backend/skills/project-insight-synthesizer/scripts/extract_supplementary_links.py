#!/usr/bin/env python3
"""Extract ess-kms knowledge-base links embedded in meeting supplementary documents (PPTX/XLSX/PDF).

This script ONLY pulls out the `https://ess-kms.edgecenter.io/{path}` links found in the
supplementary files and writes them to a manifest. It deliberately does NOT fetch the page
content: the GraphQL fetch (which needs a Bearer token) is done by the clawpm backend so the
token never has to be exposed inside the agent container.
"""

from __future__ import annotations

import argparse
import json
import re
from datetime import datetime, timezone
from pathlib import Path
from urllib.parse import urlsplit

KMS_HOST = "ess-kms.edgecenter.io"
DEFAULT_LOCALE = "en"
SUPPORTED_SUFFIXES = {".pptx", ".xlsx", ".pdf"}
# Matches bare/anchored ess-kms URLs anywhere in free text. The path charset is limited to
# RFC-3986 URL characters so it stops cleanly at CJK text / full-width punctuation that may
# immediately follow a URL with no separating whitespace.
URL_RE = re.compile(r"https?://ess-kms\.edgecenter\.io/[A-Za-z0-9\-._~:/?#\[\]@!$&'()*+,;=%]*", re.IGNORECASE)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Extract ess-kms.edgecenter.io links from PPTX/XLSX/PDF supplementary files."
    )
    parser.add_argument("inputs", type=Path, nargs="+", help="Supplementary document paths (.pptx/.xlsx/.pdf)")
    parser.add_argument("--output-dir", type=Path, required=True, help="Directory to write the manifest into")
    parser.add_argument("--manifest", type=Path, help="Destination manifest JSON path (default: <output-dir>/links-manifest.json)")
    args = parser.parse_args()

    args.output_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = (args.manifest or args.output_dir / "links-manifest.json").resolve()

    seen_paths: set[str] = set()
    links: list[dict] = []
    skipped: list[str] = []

    for input_path in args.inputs:
        if not input_path.exists():
            skipped.append(f"{input_path}: 檔案不存在")
            continue
        suffix = input_path.suffix.lower()
        if suffix not in SUPPORTED_SUFFIXES:
            skipped.append(f"{input_path}: 不支援的格式（僅支援 .pptx/.xlsx/.pdf）")
            continue
        try:
            if suffix == ".pptx":
                found = extract_pptx(input_path)
            elif suffix == ".xlsx":
                found = extract_xlsx(input_path)
            else:
                found = extract_pdf(input_path)
        except RuntimeError as exc:
            skipped.append(f"{input_path}: {exc}")
            continue

        for item in found:
            path = item["path"]
            if not path or path in seen_paths:
                continue
            seen_paths.add(path)
            links.append(item)

    manifest = {
        "generatedAt": datetime.now(timezone.utc).isoformat(),
        "links": links,
        "skipped": skipped,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Extracted {len(links)} ess-kms link(s) from {len(args.inputs)} file(s); manifest saved to {manifest_path}")
    for link in links:
        print(f"  {link['path']}  <-  {link['sourceFile']} ({link['location']})")
    for msg in skipped:
        print(f"  [skip] {msg}")


def _normalize(url: str, source_file: str, location: str) -> dict | None:
    """Turn a raw ess-kms URL into a manifest entry. Wiki.js URLs are structured as
    /{locale}/{path}, and the GraphQL singleByPath query wants the path WITHOUT the locale
    (locale passed separately), so split the leading locale segment off. Returns None if the
    host doesn't match or the path is empty."""
    url = url.strip().strip(".,;)]}>」』】）］｝、，。；：！？…\"'")
    parts = urlsplit(url)
    if parts.hostname is None or parts.hostname.lower() != KMS_HOST:
        return None
    segments = [s for s in parts.path.split("/") if s]
    if not segments:
        return None
    locale = DEFAULT_LOCALE
    # first segment like `en` / `zh` / `zh-tw` is treated as the locale
    if re.fullmatch(r"[a-z]{2}(-[a-z]{2})?", segments[0], re.IGNORECASE):
        locale = segments[0].lower()
        segments = segments[1:]
    if not segments:
        return None
    page_path = "/".join(segments)
    return {
        "url": f"{parts.scheme}://{parts.hostname}/{locale}/{page_path}",
        "path": page_path,
        "locale": locale,
        "sourceFile": source_file,
        "location": location,
    }


def _harvest_text(text: str, source_file: str, location: str) -> list[dict]:
    results = []
    for match in URL_RE.findall(text or ""):
        entry = _normalize(match, source_file, location)
        if entry:
            results.append(entry)
    return results


def extract_pptx(path: Path) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to read PPTX files (pip install python-pptx)") from exc

    results: list[dict] = []
    prs = Presentation(path)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        location = f"投影片 {slide_idx}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Explicit hyperlink on the run
                    addr = getattr(getattr(run, "hyperlink", None), "address", None)
                    if addr:
                        entry = _normalize(addr, path.name, location)
                        if entry:
                            results.append(entry)
                    # Bare URL typed as plain text
                    results.extend(_harvest_text(run.text, path.name, location))
    return results


def extract_xlsx(path: Path) -> list[dict]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to read XLSX files (pip install openpyxl)") from exc

    results: list[dict] = []
    wb = load_workbook(path, data_only=True)
    for ws in wb.worksheets:
        location = f"工作表「{ws.title}」"
        for row in ws.iter_rows():
            for cell in row:
                link = getattr(cell, "hyperlink", None)
                target = getattr(link, "target", None) if link else None
                if target:
                    entry = _normalize(target, path.name, location)
                    if entry:
                        results.append(entry)
                if isinstance(cell.value, str):
                    results.extend(_harvest_text(cell.value, path.name, location))
    return results


def extract_pdf(path: Path) -> list[dict]:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("pymupdf is required to read PDFs (pip install pymupdf)") from exc

    results: list[dict] = []
    doc = fitz.open(str(path))
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            location = f"第 {page_index + 1} 頁"
            for link in page.get_links():
                uri = link.get("uri")
                if uri:
                    entry = _normalize(uri, path.name, location)
                    if entry:
                        results.append(entry)
            results.extend(_harvest_text(page.get_text(), path.name, location))
    finally:
        doc.close()
    return results


if __name__ == "__main__":
    main()
