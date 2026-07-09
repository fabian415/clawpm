#!/usr/bin/env python3
"""
deck_builder.py — 將一份 JSON deck spec 渲染成 .pptx。

設計理念
========
不讓 LLM 直接手寫 python-pptx 座標(易變形、易溢出),而是由 agent 產出一份
結構化的 JSON deck spec(每頁的類型、標題、條列、要放哪幾張圖),再由本渲染器
負責實際排版。版面品質穩定、可重現、可除錯。

- 開啟參考模板(references/template.pptx)以「繼承主題字型 / 配色 / 16:9 尺寸」,
  但清空模板原本的 9 張範例投影片,改由本程式依 spec 動態產生 N 頁。
- 圖片鐵則:spec 內引用的每一張圖都必須在 assets-root 下真實存在,否則直接報錯,
  逼 agent 使用真實的 _assets 圖片,而非虛構。
- 圖片一律等比例縮放置入(不變形),並在框內置中。

用法
====
    python3 deck_builder.py \
        --spec deck.spec.json \
        --template <skill>/references/template.pptx \
        --assets-root /home/node/.openclaw/workspace/project-insights \
        --out output.pptx

deck spec 格式見同目錄 references/deck-spec-schema.md。
"""
import argparse
import json
import os
import random
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 主題配色(取自 template.pptx theme1.xml)──────────────────────────────────
PALETTE = {
    "dk1":     RGBColor(0x00, 0x00, 0x00),
    "lt1":     RGBColor(0xFF, 0xFF, 0xFF),
    "olive":   RGBColor(0x67, 0x6A, 0x55),   # dk2 深橄欖(主文字/標題)
    "cream":   RGBColor(0xEA, 0xEB, 0xDE),   # lt2 米色(區塊底)
    "green":   RGBColor(0x72, 0xA3, 0x76),   # accent1 主綠(強調)
    "green2":  RGBColor(0xB0, 0xCC, 0xB0),   # accent2 淺綠
    "blue":    RGBColor(0xA8, 0xCD, 0xD7),   # accent3 藍
    "sand":    RGBColor(0xC0, 0xBE, 0xAF),   # accent4 沙
    "gold":    RGBColor(0xCE, 0xC5, 0x97),   # accent5 金
    "pink":    RGBColor(0xE8, 0xB7, 0xB7),   # accent6 粉
    "red":     RGBColor(0xDB, 0x53, 0x53),   # hlink 紅(重點/警示)
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "ink":     RGBColor(0x3A, 0x3D, 0x30),   # 深一點的文字色
    "muted":   RGBColor(0x8A, 0x8C, 0x7C),   # 次要文字
    "navy":    RGBColor(0x00, 0x20, 0x60),   # 模板原始封面標題色(深藍)
}
# 章節/卡片輪播用的強調色(生動活潑)
ACCENT_CYCLE = ["green", "blue", "gold", "pink", "sand", "green2"]

FONT_EA = "微軟正黑體"     # 中文字型(模板 East-Asian font)
FONT_LATIN = "Calibri"     # 西文字型(模板 Latin font)

# ── 風格預設(讓使用者可選幾種視覺走向,但都在 ADVANTECH 模板色盤內)──────────────
# 由 meta.style 指定。primary=主題強調色、accentOrder=輪播順序、font_scale=字級倍率、
# card=中性卡片樣式(filled 米色實心+陰影 / outline 白底細框,較留白極簡)。
FONT_SCALE = 1.0
PRIMARY = "green"
CARD_STYLE = "filled"

# 版面變體種子:每份簡報一個,讓「排版方式」的預設每次不同(即使 AI 沒明確指定 variant),
# 避免每次生成的 layout 都長一樣。slide 上的 variant 欄位一律優先於這個預設。
SEED = 0
# 幾種 olive 角形章節版面(皆為橄欖色在左),輪替讓章節頁剪影每份不同。
SECTION_LAYOUTS = [67, 69, 44]


def _seed_pick(options, salt=0):
    """依 SEED + salt 從 options 挑一個(穩定但每份不同)。"""
    if not options:
        return None
    return options[(SEED + salt) % len(options)]


def _variant(s, key, options, salt=0):
    """取得某頁的版面變體:slide 明確指定優先,否則依 SEED 挑預設。"""
    v = s.get(key)
    return v if v in options else _seed_pick(options, salt)

STYLES = {
    "professional": {"primary": "green", "accentOrder": ["green", "blue", "gold", "pink", "sand"],
                     "font_scale": 1.0, "card": "filled"},
    "vivid":        {"primary": "blue",  "accentOrder": ["blue", "pink", "gold", "green", "sand"],
                     "font_scale": 1.0, "card": "filled"},
    "minimal":      {"primary": "olive", "accentOrder": ["olive", "green", "sand", "blue", "gold"],
                     "font_scale": 1.06, "card": "outline"},
    "warm":         {"primary": "gold",  "accentOrder": ["gold", "pink", "green", "sand", "blue"],
                     "font_scale": 1.0, "card": "filled"},
}


def _apply_style(meta):
    """依 meta.style 套用風格預設;meta.accentOrder(若有)一律優先於預設。"""
    global FONT_SCALE, PRIMARY, CARD_STYLE
    st = STYLES.get(meta.get("style"))
    if st:
        FONT_SCALE = st["font_scale"]
        PRIMARY = st["primary"]
        CARD_STYLE = st["card"]
        if not meta.get("accentOrder"):
            ACCENT_CYCLE[:] = st["accentOrder"]

# ── 模板 layout 對照(AI_Agent_Innovation_Contest_2026_Template.pptx / ADVANTECH）──
# 每頁建立在模板既有的 slide layout 上，以繼承企業識別(右下 ADVANTECH logo、
# 招牌橄欖綠角形色塊、母片背景),再把內容畫上去。**不畫不透明滿版底色**，
# 否則會蓋掉這些識別。以下為在此模板中經目視確認的 layout index:
LAYOUT_COVER = 0      # 品牌封面(ADVANTECH logo + Edge Computing 標語 + 天空圖)
LAYOUT_CONTENT = 1    # 乾淨白底 + 右下 logo(內容頁基底)
LAYOUT_SECTION = 67   # 橄欖綠角形色塊(章節分隔)
LAYOUT_CLOSING = 69   # 橄欖綠角形色塊(結尾)


def _layout(prs, idx):
    layouts = prs.slide_layouts
    return layouts[idx] if 0 <= idx < len(layouts) else layouts[0]

# ── 版面常數(EMU)。1 inch = 914400 EMU ────────────────────────────────────
EMU_IN = 914400
SLIDE_W = 12192000
SLIDE_H = 6858000
MARGIN = int(0.6 * EMU_IN)
CONTENT_W = SLIDE_W - 2 * MARGIN


def IN(v):
    return Emu(int(v * EMU_IN))


# ── 低階繪圖工具 ─────────────────────────────────────────────────────────────

def _delete_all_slides(prs):
    """清空模板原本的範例投影片,只留 master / layout / theme。"""
    sld_id_lst = prs.slides._sldIdLst
    rel_id_attr = qn("r:id")
    for sld_id in list(sld_id_lst):
        rId = sld_id.get(rel_id_attr)
        if rId in prs.part.rels:
            prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id)


def _new_slide(prs, layout, bg=None):
    """在指定的模板 layout 上開一張新投影片。

    移除 layout 帶進來的 placeholder(只清空提示框,母片/layout 的背景圖與 ADVANTECH
    logo 仍會保留),之後把內容自己畫上去。預設 **不畫滿版底色**(bg=None),讓模板的
    企業識別透出來;只有需要純色卡片時才由呼叫端另外畫局部形狀。
    """
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    if bg is not None:
        rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=bg, line=None)
    return slide


def _strip_slide_numbers(prs):
    """移除母片/版面中的頁碼:含 slidenum 欄位的形狀,以及 type=sldNum 的 placeholder。
    模板母片自帶頁碼欄位,會顯示在每頁右下角;使用者要求不顯示頁碼,故一律清除。"""
    def strip(container):
        try:
            shapes = container.shapes
        except Exception:
            return
        for shape in list(shapes):
            el = shape._element
            ph = el.find('.//' + qn('p:ph'))
            is_num_ph = ph is not None and ph.get('type') == 'sldNum'
            has_num_fld = any(f.get('type') == 'slidenum' for f in el.iter(qn('a:fld')))
            if is_num_ph or has_num_fld:
                el.getparent().remove(el)

    for master in prs.slide_masters:
        strip(master)
        for layout in master.slide_layouts:
            strip(layout)


# 圓角半徑(占短邊比例)。預設偏小,呈現精緻的微圓角而非陽春的大圓弧。
CORNER = 0.045


def rect(slide, left, top, width, height, fill=None, line=None,
         shape=MSO_SHAPE.RECTANGLE, line_w=None, shadow=False):
    sp = slide.shapes.add_shape(shape, Emu(int(left)), Emu(int(top)),
                                Emu(int(width)), Emu(int(height)))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = CORNER
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = PALETTE[fill] if isinstance(fill, str) else fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = PALETTE[line] if isinstance(line, str) else line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    return sp


def _soft_shadow(sp):
    """替形狀加一層柔和陰影,讓卡片浮起來(生動活潑)。"""
    spPr = sp._element.spPr
    effLst = spPr.find(qn("a:effectLst"))
    if effLst is None:
        effLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effLst)
    shdw = effLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "110000", "dist": "25400", "dir": "5400000", "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "676A55"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "13000"})
    clr.append(alpha)
    shdw.append(clr)
    effLst.append(shdw)


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT, wrap=True, autofit=True):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)),
                                  Emu(int(width)), Emu(int(height)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    # 縮放以符合大小:內容超出框就自動縮字,避免文字溢出、與相鄰區塊重疊。
    if autofit:
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
    return tb


def _card(slide, left, top, width, height):
    """中性內容卡片。filled=米色實心+陰影;outline=白底細框(極簡風,較留白)。"""
    if CARD_STYLE == "outline":
        return rect(slide, left, top, width, height, fill="white", line="sand",
                    line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    return rect(slide, left, top, width, height, fill="cream", line=None,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)


def _set_run(run, text, size, color="ink", bold=False, italic=False):
    run.text = text
    f = run.font
    f.size = Pt(round(size * FONT_SCALE, 1))
    f.bold = bold
    f.italic = italic
    f.name = FONT_LATIN
    f.color.rgb = PALETTE[color] if isinstance(color, str) else color
    # 指定 East-Asian 字型,中文才會套用微軟正黑體
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_EA)


def add_para(tb, text, size, color="ink", bold=False, italic=False,
             align=PP_ALIGN.LEFT, space_after=6, space_before=0,
             bullet=None, level=0, line_spacing=None, first=False):
    tf = tb.text_frame
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    _set_run(run, text, size, color=color, bold=bold, italic=italic)
    if bullet:
        _apply_bullet(p, bullet)
    return p


def _apply_bullet(p, char="•"):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buFont)
    pPr.append(buChar)


def add_image_fit(slide, path, box_l, box_t, box_w, box_h, center=True,
                  shadow=True, round_corners=False):
    """等比例縮放將圖片塞進指定框內並置中。回傳實際 (l,t,w,h)。"""
    pic = slide.shapes.add_picture(path, Emu(int(box_l)), Emu(int(box_t)))
    nw, nh = pic.width, pic.height
    if nw <= 0 or nh <= 0:
        return None
    scale = min(box_w / nw, box_h / nh)
    w = int(nw * scale)
    h = int(nh * scale)
    l = int(box_l + (box_w - w) / 2) if center else int(box_l)
    t = int(box_t + (box_h - h) / 2) if center else int(box_t)
    pic.left, pic.top, pic.width, pic.height = Emu(l), Emu(t), Emu(w), Emu(h)
    pic.shadow.inherit = False
    if shadow:
        _soft_shadow(pic)
    if round_corners:
        _round_pic(pic)
    return (l, t, w, h)


def _round_pic(pic, val=4500):
    """把圖片改成微圓角(val 為短邊千分比,4500=4.5%),取代預設過大的圓弧。"""
    spPr = pic._element.spPr
    for old in spPr.findall(qn("a:prstGeom")):
        spPr.remove(old)
    geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
    av = geom.makeelement(qn("a:avLst"), {})
    gd = av.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {val}"})
    av.append(gd)
    geom.append(av)
    spPr.append(geom)


def add_image_cover(slide, path, box_l, box_t, box_w, box_h):
    """裁切填滿:等比例縮放後裁掉溢出部分,讓圖片完整鋪滿整個框(不留白、不變形)。
    適合當背景/主視覺照片;若是流程圖/示意圖不要用這個(會被裁到),改用 add_image_fit。"""
    pic = slide.shapes.add_picture(path, Emu(int(box_l)), Emu(int(box_t)))
    nw, nh = pic.width, pic.height
    if nw <= 0 or nh <= 0:
        return None
    box_ratio = box_w / box_h
    img_ratio = nw / nh
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    pic.left, pic.top = Emu(int(box_l)), Emu(int(box_t))
    pic.width, pic.height = Emu(int(box_w)), Emu(int(box_h))
    pic.shadow.inherit = False
    return pic


# ── 內容取用工具 ─────────────────────────────────────────────────────────────

class ImageResolver:
    """把 spec 內的圖片參照解析成真實檔案路徑;不存在就報錯(圖片鐵則)。"""

    def __init__(self, assets_root):
        self.root = assets_root
        self.missing = []

    def resolve(self, ref):
        if not ref:
            return None
        cand = []
        if os.path.isabs(ref):
            cand.append(ref)
        cand.append(os.path.join(self.root, ref))
        cand.append(os.path.join(self.root, "_assets", os.path.basename(ref)))
        for c in cand:
            if os.path.isfile(c):
                return c
        self.missing.append(ref)
        return None


# ── 各頁型渲染 ───────────────────────────────────────────────────────────────

def slide_footer(slide, meta, page_no=None):
    # 頁尾文字放左下角。頁碼交給模板母片自己的頁碼欄位,不重複標;右下角保留給 ADVANTECH logo。
    txt = meta.get("footer", "")
    if not txt:
        return
    tb = textbox(slide, MARGIN, SLIDE_H - int(0.4 * EMU_IN),
                 int(6.0 * EMU_IN), int(0.28 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb, txt, 9, color="muted", first=True, space_after=0)


def title_band(slide, title, subtitle=None, accent=None):
    """頁面上緣的標題區:一條強調色豎條 + 標題 + 副標。回傳內容起始 top。"""
    accent = accent or PRIMARY
    bar_t = MARGIN
    rect(slide, MARGIN, bar_t + int(0.03 * EMU_IN), int(0.11 * EMU_IN),
         int(0.72 * EMU_IN), fill=accent, line=None)
    tb = textbox(slide, MARGIN + int(0.32 * EMU_IN), bar_t,
                 CONTENT_W - int(0.32 * EMU_IN), int(0.95 * EMU_IN))
    add_para(tb, title, 30, color="olive", bold=True, first=True, space_after=3)
    if subtitle:
        add_para(tb, subtitle, 15, color="muted", space_after=0)
    return bar_t + int(1.25 * EMU_IN)


def render_cover(prs, layout, meta, s, resolver):
    # 一律沿用模板「原始封面」(L0:ADVANTECH logo + 天空 + Edge Computing 標語),
    # 忠實重現原稿結構:Project Title / Domain—Key Area & Application / Core Team / Date。
    # 不做其他變體、不加任何賽事字樣;標題採原稿深藍色。
    slide = _new_slide(prs, _layout(prs, LAYOUT_COVER))
    title = s.get("title", meta.get("title", ""))
    domain = s.get("domain", "")
    team = s.get("team") or meta.get("team") or ""
    date = s.get("date") or meta.get("date") or ""
    L = int(0.6 * EMU_IN)

    tb = textbox(slide, L, int(2.35 * EMU_IN), int(8.6 * EMU_IN), int(1.7 * EMU_IN))
    add_para(tb, "Project Title", 15, color="muted", first=True, space_after=4)
    add_para(tb, title, 30, color="navy", bold=True, space_after=0, line_spacing=1.08)

    db = textbox(slide, L, int(4.35 * EMU_IN), int(8.6 * EMU_IN), int(1.3 * EMU_IN))
    add_para(db, "Domain — Key Area & Application", 15, color="muted", first=True,
             space_after=4)
    if domain:
        add_para(db, domain, 24, color="navy", bold=True, space_after=0)

    mb = textbox(slide, L, SLIDE_H - int(1.35 * EMU_IN), int(9.0 * EMU_IN),
                 int(1.05 * EMU_IN))
    add_para(mb, f"Core Team：{team}", 15, color="olive", first=True, space_after=5)
    add_para(mb, f"Date：{date}", 15, color="olive", space_after=0)
    return slide


def render_toc(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", "目錄 / Agenda"),
                     s.get("subtitle"), accent=PRIMARY)
    items = s.get("items", [])
    n = len(items)
    col_h = SLIDE_H - top - int(0.8 * EMU_IN)
    per_col = (n + 1) // 2 if n > 5 else n
    col_w = (CONTENT_W - int(0.6 * EMU_IN)) // 2 if n > 5 else CONTENT_W
    for i, item in enumerate(items):
        col = 0 if i < per_col else 1
        row = i if i < per_col else i - per_col
        rh = col_h / max(per_col, 1)
        l = MARGIN + col * (col_w + int(0.6 * EMU_IN))
        t = top + int(row * rh)
        accent = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        # 數字徽章
        badge = int(0.58 * EMU_IN)
        rect(slide, l, t + int(0.08 * EMU_IN), badge, badge, fill=accent,
             line=None, shape=MSO_SHAPE.OVAL)
        nb = textbox(slide, l, t + int(0.08 * EMU_IN), badge, badge,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_para(nb, str(i + 1).zfill(2), 17, color="white", bold=True,
                 first=True, align=PP_ALIGN.CENTER, space_after=0)
        it = textbox(slide, l + badge + int(0.26 * EMU_IN), t,
                     col_w - badge - int(0.36 * EMU_IN), int(rh),
                     anchor=MSO_ANCHOR.MIDDLE)
        title = item if isinstance(item, str) else item.get("title", "")
        add_para(it, title, 18, color="olive", bold=True, first=True, space_after=2)
        if isinstance(item, dict) and item.get("desc"):
            add_para(it, item["desc"], 12, color="muted", space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


def render_section(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    # 建立在模板橄欖綠角形 layout(L67)上,不畫滿版底色。標題白字放左側橄欖色塊區。
    # 橄欖色塊是斜角(中央最寬、上下較窄),文字框必須夠窄並置中對齊,長標題才不會溢出到白底。
    img = resolver.resolve(s.get("image"))
    if img:
        # 圖片放大填滿右側白區(白區約 x=6.7"~13.3"),置中、加圓角與陰影,避免小圖漂浮。
        add_image_fit(slide, img, SLIDE_W - int(6.7 * EMU_IN), int(1.15 * EMU_IN),
                      int(6.0 * EMU_IN), int(4.7 * EMU_IN), round_corners=True)
        if s.get("caption"):
            cb = textbox(slide, SLIDE_W - int(6.7 * EMU_IN), SLIDE_H - int(1.05 * EMU_IN),
                         int(6.0 * EMU_IN), int(0.4 * EMU_IN), align=PP_ALIGN.CENTER)
            add_para(cb, s["caption"], 12, color="muted", first=True,
                     align=PP_ALIGN.CENTER, space_after=0)
    else:
        # 無圖時,在白區放一個大型淡色章節編號當視覺重心,避免右側空洞。
        num = _leading_number(s.get("kicker"))
        if num:
            nb = textbox(slide, SLIDE_W - int(5.2 * EMU_IN), int(1.6 * EMU_IN),
                         int(4.4 * EMU_IN), int(3.6 * EMU_IN),
                         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            add_para(nb, num, 200, color="cream", bold=True, first=True,
                     align=PP_ALIGN.CENTER, space_after=0)
    tb = textbox(slide, int(0.7 * EMU_IN), int(2.3 * EMU_IN),
                 int(3.6 * EMU_IN), int(2.7 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
    if s.get("kicker"):
        add_para(tb, s["kicker"], 15, color="cream", bold=True, first=True,
                 space_after=10)
    add_para(tb, s.get("title", ""), 30, color="white", bold=True,
             space_after=8, line_spacing=1.06,
             first=not s.get("kicker"))
    if s.get("subtitle"):
        add_para(tb, s["subtitle"], 14, color="cream", space_after=0,
                 line_spacing=1.12)
    return slide


def _leading_number(text):
    """從 kicker(如 '01 / Market Problem')取出開頭的數字,給無圖章節頁當大型視覺編號。"""
    import re
    m = re.match(r"\s*(\d{1,2})", str(text or ""))
    return m.group(1).zfill(2) if m else None


def render_bullets(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", ""), s.get("subtitle"),
                     accent=s.get("accent", PRIMARY))
    acc = s.get("accent", PRIMARY)
    img = resolver.resolve(s.get("image"))
    bullets = s.get("bullets", [])
    area_h = SLIDE_H - top - int(0.8 * EMU_IN)

    # 版面變體:有圖時圖片放左或右;無圖時可用卡片或條列。每份預設依 SEED 不同。
    variant = _variant(s, "variant", ["list", "cards"], salt=2)
    if img:
        side = _variant(s, "imageSide", ["right", "left"], salt=3)
        img_w = int(CONTENT_W * 0.42)
        if side == "left":
            img_l = MARGIN
            text_l = MARGIN + img_w + int(0.5 * EMU_IN)
        else:
            text_l = MARGIN
            img_l = SLIDE_W - MARGIN - img_w
        add_image_fit(slide, img, img_l, top, img_w, area_h, round_corners=True)
        if s.get("caption"):
            cb = textbox(slide, img_l, SLIDE_H - int(0.72 * EMU_IN), img_w,
                         int(0.32 * EMU_IN), align=PP_ALIGN.CENTER)
            add_para(cb, s["caption"], 11, color="muted", first=True,
                     align=PP_ALIGN.CENTER, space_after=0)
        text_w = CONTENT_W - img_w - int(0.5 * EMU_IN)
        _render_bullet_list(slide, bullets, text_l, top, text_w, area_h, accent=acc)
    elif variant == "cards":
        _render_bullet_cards(slide, bullets, MARGIN, top, CONTENT_W, area_h, accent=acc)
    else:
        _render_bullet_list(slide, bullets, MARGIN, top, CONTENT_W, area_h, accent=acc)
    slide_footer(slide, meta, page_no)
    return slide


def _render_bullet_cards(slide, bullets, l, t, w, h, accent=PRIMARY):
    """把每個要點做成一張卡片(≤3 張並排,4 張 2×2),比純條列更有版面感。"""
    if not bullets:
        return
    n = len(bullets)
    cols = n if n <= 3 else 2
    rows = (n + cols - 1) // cols
    gap = int(0.35 * EMU_IN)
    cw = (w - (cols - 1) * gap) // cols
    ch = (h - (rows - 1) * gap) // rows
    for i, b in enumerate(bullets):
        r, c = divmod(i, cols)
        cl = l + c * (cw + gap)
        ct = t + r * (ch + gap)
        acc = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        _card(slide, cl, ct, cw, ch)
        rect(slide, cl, ct, cw, int(0.14 * EMU_IN), fill=acc, line=None)
        tb = textbox(slide, cl + int(0.28 * EMU_IN), ct + int(0.3 * EMU_IN),
                     cw - int(0.56 * EMU_IN), ch - int(0.5 * EMU_IN))
        if isinstance(b, str):
            add_para(tb, b, 16, color="ink", first=True, space_after=0, line_spacing=1.14)
        else:
            add_para(tb, b.get("heading", ""), 18, color="olive", bold=True,
                     first=True, space_after=5)
            if b.get("text"):
                add_para(tb, b["text"], 14, color="ink", space_after=0, line_spacing=1.14)


def _render_bullet_list(slide, bullets, l, t, w, h, accent=PRIMARY):
    if not bullets:
        return
    n = len(bullets)
    rh = h / n
    for i, b in enumerate(bullets):
        by = t + int(i * rh)
        acc = accent if accent else ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        # 小色點(垂直對齊第一行文字)
        dot = int(0.19 * EMU_IN)
        rect(slide, l, by + int(0.14 * EMU_IN), dot, dot, fill=acc, line=None,
             shape=MSO_SHAPE.OVAL)
        tb = textbox(slide, l + dot + int(0.26 * EMU_IN), by,
                     w - dot - int(0.34 * EMU_IN), int(rh), anchor=MSO_ANCHOR.TOP)
        if isinstance(b, str):
            add_para(tb, b, 17, color="ink", first=True, space_after=3,
                     line_spacing=1.12)
        else:
            add_para(tb, b.get("heading", ""), 19, color="olive", bold=True,
                     first=True, space_after=3)
            if b.get("text"):
                add_para(tb, b["text"], 15, color="ink", space_after=0,
                         line_spacing=1.12)


def render_two_column(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", ""), s.get("subtitle"),
                     accent=s.get("accent", PRIMARY))
    gap = int(0.5 * EMU_IN)
    col_w = (CONTENT_W - gap) // 2
    col_h = SLIDE_H - top - int(0.8 * EMU_IN)
    a0, a1 = ACCENT_CYCLE[0], ACCENT_CYCLE[1]
    variant = _variant(s, "variant", ["panels", "divider"], salt=4)
    for idx, key, acc in [(0, "left", a0), (1, "right", a1)]:
        col = s.get(key, {})
        l = MARGIN + idx * (col_w + gap)
        if variant == "divider":
            # 無卡片:色塊標頭條 + 條列,中間一條分隔線
            rect(slide, l, top, col_w, int(0.72 * EMU_IN), fill=acc, line=None,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            hb = textbox(slide, l + int(0.3 * EMU_IN), top, col_w - int(0.5 * EMU_IN),
                         int(0.72 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
            add_para(hb, col.get("heading", ""), 18, color="white", bold=True,
                     first=True, space_after=0)
            _render_bullet_list(slide, col.get("bullets", []),
                                l + int(0.1 * EMU_IN), top + int(1.0 * EMU_IN),
                                col_w - int(0.3 * EMU_IN), col_h - int(1.3 * EMU_IN),
                                accent=acc)
        else:
            _card(slide, l, top, col_w, col_h)
            rect(slide, l, top, col_w, int(0.74 * EMU_IN), fill=acc, line=None,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            hb = textbox(slide, l + int(0.35 * EMU_IN), top, col_w - int(0.6 * EMU_IN),
                         int(0.74 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
            add_para(hb, col.get("heading", ""), 18, color="white", bold=True,
                     first=True, space_after=0)
            _render_bullet_list(slide, col.get("bullets", []),
                                l + int(0.38 * EMU_IN), top + int(1.05 * EMU_IN),
                                col_w - int(0.7 * EMU_IN),
                                col_h - int(1.35 * EMU_IN), accent=acc)
    if variant == "divider":
        rect(slide, MARGIN + col_w + gap // 2 - int(0.01 * EMU_IN),
             top + int(0.9 * EMU_IN), int(0.02 * EMU_IN), col_h - int(1.0 * EMU_IN),
             fill="sand", line=None)
    slide_footer(slide, meta, page_no)
    return slide


def render_image_full(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", ""), s.get("subtitle"),
                     accent=s.get("accent", PRIMARY))
    img = resolver.resolve(s.get("image"))
    cap_h = int(0.5 * EMU_IN) if s.get("caption") else 0
    box_h = SLIDE_H - top - int(0.7 * EMU_IN) - cap_h
    if img:
        add_image_fit(slide, img, MARGIN, top, CONTENT_W, box_h,
                      round_corners=True)
    else:
        rect(slide, MARGIN, top, CONTENT_W, box_h, fill="cream", line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if s.get("caption"):
        cb = textbox(slide, MARGIN, top + box_h + int(0.1 * EMU_IN), CONTENT_W,
                     cap_h, align=PP_ALIGN.CENTER)
        add_para(cb, s["caption"], 13, color="muted", first=True,
                 align=PP_ALIGN.CENTER, space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


def render_image_grid(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", ""), s.get("subtitle"),
                     accent=s.get("accent", PRIMARY))
    images = [im for im in s.get("images", []) if resolver.resolve(
        im.get("path") if isinstance(im, dict) else im)]
    n = len(images)
    if n == 0:
        slide_footer(slide, meta, page_no)
        return slide
    cols = 1 if n == 1 else (2 if n <= 4 else 3)
    rows = (n + cols - 1) // cols
    gap = int(0.35 * EMU_IN)
    area_h = SLIDE_H - top - int(0.8 * EMU_IN)
    cell_w = (CONTENT_W - (cols - 1) * gap) // cols
    cell_h = (area_h - (rows - 1) * gap) // rows
    for i, im in enumerate(images):
        r, c = divmod(i, cols)
        l = MARGIN + c * (cell_w + gap)
        t = top + r * (cell_h + gap)
        path = resolver.resolve(im.get("path") if isinstance(im, dict) else im)
        cap = im.get("caption") if isinstance(im, dict) else None
        cap_h = int(0.34 * EMU_IN) if cap else 0
        _card(slide, l, t, cell_w, cell_h)
        add_image_fit(slide, path, l + int(0.08 * EMU_IN), t + int(0.08 * EMU_IN),
                      cell_w - int(0.16 * EMU_IN),
                      cell_h - int(0.16 * EMU_IN) - cap_h,
                      shadow=False, round_corners=True)
        if cap:
            cb = textbox(slide, l, t + cell_h - cap_h, cell_w, cap_h,
                         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            add_para(cb, cap, 12, color="olive", first=True,
                     align=PP_ALIGN.CENTER, space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


def render_stats(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", ""), s.get("subtitle"),
                     accent=s.get("accent", PRIMARY))
    stats = s.get("stats", [])
    n = max(len(stats), 1)
    area_h = SLIDE_H - top - int(0.8 * EMU_IN)
    gap = int(0.45 * EMU_IN)
    # 變體:row(單排,大字最醒目)/ grid(2 欄,4 個時較平衡)。>4 個一律換行。
    variant = _variant(s, "variant", ["row", "grid"], salt=5)
    cols = n if (variant == "row" and n <= 4) else (2 if n <= 4 else 3)
    rows = (n + cols - 1) // cols
    card_w = (CONTENT_W - (cols - 1) * gap) // cols
    card_h = int(3.0 * EMU_IN) if rows == 1 else int((area_h - (rows - 1) * gap) / rows)
    t0 = top + max(0, (area_h - (card_h * rows + (rows - 1) * gap)) // 2)
    for i, st in enumerate(stats):
        acc = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        r, c = divmod(i, cols)
        l = MARGIN + c * (card_w + gap)
        t = t0 + r * (card_h + gap)
        _card(slide, l, t, card_w, card_h)
        rect(slide, l, t, card_w, int(0.16 * EMU_IN), fill=acc, line=None)
        vh = int(card_h * 0.5)
        vb = textbox(slide, l, t + int(0.35 * EMU_IN), card_w, vh,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_para(vb, str(st.get("value", "")), 46 if rows == 1 else 38, color=acc,
                 bold=True, first=True, align=PP_ALIGN.CENTER, space_after=0)
        lb = textbox(slide, l + int(0.25 * EMU_IN), t + int(0.4 * EMU_IN) + vh,
                     card_w - int(0.5 * EMU_IN), card_h - vh - int(0.5 * EMU_IN),
                     anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER)
        add_para(lb, str(st.get("label", "")), 16, color="olive", bold=True,
                 first=True, align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.1)
    if s.get("note"):
        nb = textbox(slide, MARGIN, SLIDE_H - int(0.7 * EMU_IN), CONTENT_W,
                     int(0.4 * EMU_IN))
        add_para(nb, s["note"], 12, color="muted", first=True, space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


def render_summary(prs, layout, meta, s, resolver, page_no):
    """Executive-summary 表格式:左標籤 + 右內容,一列一項。"""
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", "Executive Summary"),
                     s.get("subtitle"), accent=PRIMARY)
    rows = s.get("rows", [])
    n = max(len(rows), 1)
    area_h = SLIDE_H - top - int(0.8 * EMU_IN)
    rh = min(area_h / n, int(1.15 * EMU_IN))
    label_w = int(3.0 * EMU_IN)
    # 若各列合計未填滿內容區,整組垂直置中
    y0 = top + max(0, (area_h - int(rh * n)) // 2)
    for i, row in enumerate(rows):
        acc = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        t = y0 + int(i * rh)
        rect(slide, MARGIN, t + int(0.06 * EMU_IN), label_w,
             int(rh) - int(0.16 * EMU_IN), fill=acc, line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        lb = textbox(slide, MARGIN + int(0.25 * EMU_IN), t,
                     label_w - int(0.4 * EMU_IN), int(rh),
                     anchor=MSO_ANCHOR.MIDDLE)
        add_para(lb, row.get("label", ""), 16, color="white", bold=True,
                 first=True, space_after=0)
        cb = textbox(slide, MARGIN + label_w + int(0.4 * EMU_IN), t,
                     CONTENT_W - label_w - int(0.4 * EMU_IN), int(rh),
                     anchor=MSO_ANCHOR.MIDDLE)
        add_para(cb, row.get("text", ""), 15, color="ink", first=True,
                 space_after=0, line_spacing=1.12)
    slide_footer(slide, meta, page_no)
    return slide


def render_roadmap(prs, layout, meta, s, resolver, page_no):
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", "Roadmap"), s.get("subtitle"),
                     accent=PRIMARY)
    phases = s.get("phases", [])
    n = max(len(phases), 1)
    gap = int(0.35 * EMU_IN)
    card_w = (CONTENT_W - (n - 1) * gap) // n
    t = top + int(0.5 * EMU_IN)
    card_h = int(3.2 * EMU_IN)
    # 連接時間軸
    rect(slide, MARGIN, t + int(0.32 * EMU_IN), CONTENT_W, int(0.05 * EMU_IN),
         fill="green2", line=None)
    for i, ph in enumerate(phases):
        acc = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        l = MARGIN + i * (card_w + gap)
        node = int(0.34 * EMU_IN)
        rect(slide, l + card_w // 2 - node // 2, t + int(0.18 * EMU_IN),
             node, node, fill=acc, line="white", line_w=2, shape=MSO_SHAPE.OVAL)
        ct = t + int(0.8 * EMU_IN)
        _card(slide, l, ct, card_w, card_h)
        pb = textbox(slide, l + int(0.25 * EMU_IN), ct + int(0.22 * EMU_IN),
                     card_w - int(0.5 * EMU_IN), int(0.55 * EMU_IN))
        add_para(pb, ph.get("period", ""), 18, color=acc, bold=True, first=True,
                 space_after=0)
        bb = textbox(slide, l + int(0.25 * EMU_IN), ct + int(0.85 * EMU_IN),
                     card_w - int(0.5 * EMU_IN), card_h - int(1.0 * EMU_IN))
        items = ph.get("items") or ([ph.get("text")] if ph.get("text") else [])
        for j, it in enumerate(items):
            add_para(bb, it, 14, color="ink", first=(j == 0), space_after=6,
                     bullet="•", line_spacing=1.12)
    slide_footer(slide, meta, page_no)
    return slide


def render_closing(prs, layout, meta, s, resolver, page_no):
    # 建立在模板橄欖綠角形 layout(L69)上,標題白字放左側橄欖區。
    slide = _new_slide(prs, layout)
    tb = textbox(slide, int(0.8 * EMU_IN), int(2.4 * EMU_IN),
                 int(4.7 * EMU_IN), int(2.2 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb, s.get("title", "Thank You"), 44, color="white", bold=True,
             first=True, space_after=10)
    if s.get("subtitle"):
        add_para(tb, s["subtitle"], 16, color="cream", space_after=0,
                 line_spacing=1.15)
    return slide


def render_quote(prs, layout, meta, s, resolver, page_no):
    """金句頁:大引號 + 置中粗體引文 + 出處。用於故事轉折/核心主張,增加節奏變化。"""
    slide = _new_slide(prs, layout)
    acc = s.get("accent", PRIMARY)
    qb = textbox(slide, MARGIN, int(1.0 * EMU_IN), int(1.8 * EMU_IN), int(1.6 * EMU_IN))
    add_para(qb, "“", 130, color=acc, bold=True, first=True, space_after=0)
    tb = textbox(slide, int(1.7 * EMU_IN), int(2.0 * EMU_IN),
                 SLIDE_W - int(3.4 * EMU_IN), int(3.0 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb, s.get("text", ""), 30, color="olive", bold=True, first=True,
             line_spacing=1.22, space_after=12)
    if s.get("attribution"):
        add_para(tb, "— " + s["attribution"], 16, color="muted", space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


def render_hero_image(prs, layout, meta, s, resolver, page_no):
    """滿版主圖頁:右側大圖鋪滿、左側色塊面板放標題與短文。用於畫面感強的關鍵頁。"""
    slide = _new_slide(prs, layout)
    acc = s.get("accent", PRIMARY)
    panel_w = int(4.9 * EMU_IN)
    img = resolver.resolve(s.get("image"))
    if img:
        add_image_cover(slide, img, panel_w - int(0.1 * EMU_IN), 0,
                        SLIDE_W - panel_w + int(0.1 * EMU_IN), SLIDE_H)
    rect(slide, 0, 0, panel_w, SLIDE_H, fill=acc, line=None)
    tb = textbox(slide, int(0.65 * EMU_IN), int(1.9 * EMU_IN),
                 panel_w - int(1.1 * EMU_IN), int(3.4 * EMU_IN), anchor=MSO_ANCHOR.MIDDLE)
    if s.get("kicker"):
        add_para(tb, s["kicker"], 14, color="cream", bold=True, first=True, space_after=10)
    add_para(tb, s.get("title", ""), 30, color="white", bold=True,
             first=not s.get("kicker"), space_after=10, line_spacing=1.08)
    if s.get("text"):
        add_para(tb, s["text"], 15, color="cream", space_after=0, line_spacing=1.18)
    return slide


def render_big_number(prs, layout, meta, s, resolver, page_no):
    """單一巨大數據頁:一個關鍵數字撐滿版面,底下 label 與說明。強調最有力的一個證據。"""
    slide = _new_slide(prs, layout)
    acc = s.get("accent", PRIMARY)
    rect(slide, SLIDE_W // 2 - int(0.6 * EMU_IN), int(1.5 * EMU_IN),
         int(1.2 * EMU_IN), int(0.12 * EMU_IN), fill=acc, line=None)
    vb = textbox(slide, 0, int(1.9 * EMU_IN), SLIDE_W, int(2.3 * EMU_IN),
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_para(vb, str(s.get("value", "")), 130, color=acc, bold=True, first=True,
             align=PP_ALIGN.CENTER, space_after=0)
    lb = textbox(slide, int(1.5 * EMU_IN), int(4.35 * EMU_IN),
                 SLIDE_W - int(3.0 * EMU_IN), int(0.8 * EMU_IN), align=PP_ALIGN.CENTER)
    add_para(lb, str(s.get("label", "")), 22, color="olive", bold=True, first=True,
             align=PP_ALIGN.CENTER, space_after=4)
    if s.get("context"):
        cb = textbox(slide, int(2.0 * EMU_IN), int(5.2 * EMU_IN),
                     SLIDE_W - int(4.0 * EMU_IN), int(1.0 * EMU_IN), align=PP_ALIGN.CENTER)
        add_para(cb, s["context"], 15, color="muted", first=True,
                 align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.15)
    slide_footer(slide, meta, page_no)
    return slide


def render_process(prs, layout, meta, s, resolver, page_no):
    """流程步驟頁:水平 N 個步驟卡 + 之間箭頭。用於工作流程/導入路徑,取代死板條列。"""
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", ""), s.get("subtitle"),
                     accent=s.get("accent", PRIMARY))
    steps = s.get("steps", [])
    n = max(len(steps), 1)
    arrow = int(0.5 * EMU_IN)
    card_w = (CONTENT_W - (n - 1) * arrow) // n
    area_h = SLIDE_H - top - int(0.8 * EMU_IN)
    card_h = int(3.0 * EMU_IN)
    t = top + max(int(0.1 * EMU_IN), (area_h - card_h) // 2)
    for i, st in enumerate(steps):
        acc = ACCENT_CYCLE[i % len(ACCENT_CYCLE)]
        l = MARGIN + i * (card_w + arrow)
        _card(slide, l, t, card_w, card_h)
        badge = int(0.66 * EMU_IN)
        rect(slide, l + card_w // 2 - badge // 2, t + int(0.35 * EMU_IN), badge, badge,
             fill=acc, line=None, shape=MSO_SHAPE.OVAL)
        nb = textbox(slide, l + card_w // 2 - badge // 2, t + int(0.35 * EMU_IN),
                     badge, badge, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_para(nb, str(i + 1), 20, color="white", bold=True, first=True,
                 align=PP_ALIGN.CENTER, space_after=0)
        tb = textbox(slide, l + int(0.25 * EMU_IN), t + int(1.25 * EMU_IN),
                     card_w - int(0.5 * EMU_IN), card_h - int(1.4 * EMU_IN),
                     anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER)
        title = st if isinstance(st, str) else st.get("title", "")
        add_para(tb, title, 17, color="olive", bold=True, first=True,
                 align=PP_ALIGN.CENTER, space_after=4)
        if isinstance(st, dict) and st.get("text"):
            add_para(tb, st["text"], 13, color="ink", align=PP_ALIGN.CENTER,
                     space_after=0, line_spacing=1.12)
        # 步驟間箭頭
        if i < n - 1:
            ab = textbox(slide, l + card_w, t, arrow, card_h,
                         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            add_para(ab, "→", 24, color="green2", bold=True, first=True,
                     align=PP_ALIGN.CENTER, space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


def render_sources(prs, layout, meta, s, resolver, page_no):
    """市場洞察來源頁:列出外部研究引用(標題 + 連結),對應「引用附連結」的要求。"""
    slide = _new_slide(prs, layout)
    top = title_band(slide, s.get("title", "市場洞察來源 / Sources"),
                     s.get("subtitle"), accent=PRIMARY)
    items = s.get("sources", [])
    n = max(len(items), 1)
    area_h = SLIDE_H - top - int(0.8 * EMU_IN)
    rh = min(area_h / n, int(0.85 * EMU_IN))
    for i, it in enumerate(items):
        t = top + int(i * rh)
        title = it if isinstance(it, str) else it.get("title", "")
        url = it.get("url", "") if isinstance(it, dict) else ""
        dot = int(0.16 * EMU_IN)
        rect(slide, MARGIN, t + int(0.12 * EMU_IN), dot, dot,
             fill=ACCENT_CYCLE[i % len(ACCENT_CYCLE)], line=None, shape=MSO_SHAPE.OVAL)
        tb = textbox(slide, MARGIN + dot + int(0.24 * EMU_IN), t,
                     CONTENT_W - dot - int(0.3 * EMU_IN), int(rh))
        add_para(tb, title, 15, color="olive", bold=True, first=True, space_after=2)
        if url:
            add_para(tb, url, 12, color="blue", space_after=0)
    slide_footer(slide, meta, page_no)
    return slide


RENDERERS = {
    "cover": render_cover,
    "toc": render_toc,
    "section": render_section,
    "bullets": render_bullets,
    "two_column": render_two_column,
    "quote": render_quote,
    "hero_image": render_hero_image,
    "big_number": render_big_number,
    "process": render_process,
    "sources": render_sources,
    "image_full": render_image_full,
    "image_grid": render_image_grid,
    "stats": render_stats,
    "summary": render_summary,
    "roadmap": render_roadmap,
    "closing": render_closing,
}


def build(spec, template, assets_root, out):
    prs = Presentation(template)
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    _delete_all_slides(prs)
    _strip_slide_numbers(prs)   # 使用者要求不顯示頁碼
    resolver = ImageResolver(assets_root)
    meta = spec.get("meta", {})

    # 風格預設(使用者可選)+ 版面變化引擎(每次可帶不同 accentOrder / layoutSeed,避免每份都一樣)。
    _apply_style(meta)
    global SEED
    seed = meta.get("layoutSeed")
    SEED = int(seed) if isinstance(seed, (int, float)) else random.randint(0, 9999)
    order = meta.get("accentOrder")
    if isinstance(order, list) and len(order) >= 3 and all(o in PALETTE for o in order):
        ACCENT_CYCLE[:] = order

    # 每種頁型建立在對應的模板 layout 上,以繼承 ADVANTECH 企業識別
    layout_for = {"cover": LAYOUT_COVER, "section": LAYOUT_SECTION,
                  "closing": LAYOUT_CLOSING}

    page_no = 0
    section_i = 0
    for idx, s in enumerate(spec.get("slides", [])):
        stype = s.get("type", "bullets")
        renderer = RENDERERS.get(stype)
        if renderer is None:
            sys.stderr.write(f"[warn] 未知頁型 '{stype}',以 bullets 取代\n")
            renderer = render_bullets
            stype = "bullets"
        if stype == "section":
            # 章節頁在數種橄欖角形版面間輪替,讓剪影每份/每章不同
            layout = _layout(prs, _seed_pick(SECTION_LAYOUTS, salt=section_i))
            section_i += 1
        else:
            layout = _layout(prs, layout_for.get(stype, LAYOUT_CONTENT))
        # cover / section / closing 不編頁碼
        if stype in ("cover", "section", "closing"):
            if stype == "cover":
                renderer(prs, layout, meta, s, resolver)
            else:
                renderer(prs, layout, meta, s, resolver, None)
        else:
            page_no += 1
            renderer(prs, layout, meta, s, resolver, page_no)

    os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
    prs.save(out)

    result = {
        "ok": True,
        "out": out,
        "slide_count": len(prs.slides._sldIdLst),
        "missing_images": resolver.missing,
    }
    if resolver.missing:
        result["ok"] = False
        result["error"] = (
            "以下 spec 引用的圖片在 assets-root 找不到,請改用真實存在的 _assets 圖片:"
            + ", ".join(resolver.missing)
        )
    return result


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--assets-root", required=True,
                    help="圖片相對路徑(如 _assets/x.png)的解析根目錄,通常為 project-insights/")
    ap.add_argument("--out", required=True)
    ap.add_argument("--allow-missing", action="store_true",
                    help="缺圖時仍輸出(預設缺圖視為失敗)")
    args = ap.parse_args()

    with open(args.spec, "r", encoding="utf-8") as f:
        spec = json.load(f)

    result = build(spec, args.template, args.assets_root, args.out)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if result.get("missing_images") and not args.allow_missing:
        sys.exit(3)


if __name__ == "__main__":
    main()
